import importlib.metadata
import re
import struct
import zlib

import pytest


def _png_chunk(kind, data):
    return struct.pack(">I", len(data)) + kind + data + struct.pack(">I", zlib.crc32(kind + data))


def _write_tiny_png(path):
    signature = b"\x89PNG\r\n\x1a\n"
    header = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    pixels = zlib.compress(b"\x00\x20\x80\xe0")
    path.write_bytes(
        signature
        + _png_chunk(b"IHDR", header)
        + _png_chunk(b"IDAT", pixels)
        + _png_chunk(b"IEND", b"")
    )


def test_python_pptx_getting_started_examples_run_with_global_revision_refetches(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    _write_tiny_png(tmp_path / "monty-truth.png")

    # Hello World
    from rpptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle = prs.slides[0].placeholders[1]
    subtitle.text = "python-pptx was here!"

    prs.save("test.pptx")

    # Bullet slide
    from rpptx import Presentation

    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Adding a Bullet Slide"

    body_shape = prs.slides[0].shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Find the bullet slide layout"

    tf = prs.slides[0].shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Use _TextFrame.text for first bullet"
    p = prs.slides[0].shapes.placeholders[1].text_frame.paragraphs[1]
    p.level = 1

    tf = prs.slides[0].shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Use _TextFrame.add_paragraph() for subsequent bullets"
    p = prs.slides[0].shapes.placeholders[1].text_frame.paragraphs[2]
    p.level = 2

    prs.save("test.pptx")

    # add_textbox()
    from rpptx import Presentation
    from rpptx.util import Inches, Pt

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = "This is text inside a textbox"

    tf = prs.slides[0].shapes[-1].text_frame
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p = prs.slides[0].shapes[-1].text_frame.paragraphs[1]
    p.font.bold = True

    tf = prs.slides[0].shapes[-1].text_frame
    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p = prs.slides[0].shapes[-1].text_frame.paragraphs[2]
    p.font.size = Pt(40)

    prs.save("test.pptx")

    # add_picture()
    from rpptx import Presentation
    from rpptx.util import Inches

    img_path = "monty-truth.png"

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)

    slide = prs.slides[0]
    left = Inches(5)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save("test.pptx")

    # add_shape()
    from rpptx import Presentation
    from rpptx.enum.shapes import MSO_SHAPE
    from rpptx.util import Inches

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = "Adding an AutoShape"
    shapes = prs.slides[0].shapes

    left = Inches(0.93)
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = "Step 1"

    left = left + width - Inches(0.4)
    width = Inches(2.0)

    for n in range(2, 6):
        shapes = prs.slides[0].shapes
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = "Step %d" % n
        left = left + width - Inches(0.4)

    prs.save("test.pptx")

    # add_table()
    from rpptx import Presentation
    from rpptx.util import Inches

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = "Adding a Table"
    shapes = prs.slides[0].shapes

    rows = cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    table.cell(0, 0).text = "Foo"
    table.cell(0, 1).text = "Bar"

    table.cell(1, 0).text = "Baz"
    table.cell(1, 1).text = "Qux"

    prs.save("test.pptx")

    # Extract all text from slides
    from rpptx import Presentation

    path_to_presentation = "test.pptx"
    prs = Presentation(path_to_presentation)

    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    assert text_runs == ["Adding a Table"]
    assert (tmp_path / "test.pptx").is_file()


def test_lazy_collections_and_stale_handles_are_loud():
    import rpptx

    prs = rpptx.Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[6])
    held = first.shapes.add_textbox(0, 0, 100, 100)
    prs.slides.add_slide(prs.slide_layouts[6])

    assert len(prs.slides) == 2
    assert list(prs.slides[:1])[0].shapes[-1].text == ""
    with pytest.raises(rpptx.StaleElementError, match=r"revision 2.*revision 3"):
        _ = held.text


def test_omitted_placeholder_index_resolves_as_default_zero():
    import rpptx

    prs = rpptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "default zero"
    assert prs.slides[0].shapes.title.text == "default zero"


def _assert_stale_after_exactly_one_bump(rpptx, operation):
    with pytest.raises(rpptx.StaleElementError) as raised:
        operation()
    revisions = re.search(
        r"revision (\d+).*revision (\d+)", str(raised.value)
    )
    assert revisions is not None
    captured, current = map(int, revisions.groups())
    assert current == captured + 1


def _assert_exact_stale(rpptx, operation, kind, captured, current, recovery):
    expected = (
        f"{kind} handle was created at document revision {captured}, but the "
        f"document is now at revision {current} (a structural change "
        f"invalidated it). {recovery}"
    )
    with pytest.raises(rpptx.StaleElementError) as raised:
        operation()
    assert str(raised.value) == expected


def test_structural_append_invalidates_every_preexisting_path_handle():
    import rpptx

    prs = rpptx.Presentation()
    layouts = prs.slide_layouts
    layout = layouts[6]
    slides = prs.slides
    slide = slides.add_slide(layout)

    _assert_exact_stale(
        rpptx,
        lambda: len(slides),
        "slide collection",
        0,
        1,
        "Re-fetch it with prs.slides.",
    )
    _assert_exact_stale(
        rpptx,
        lambda: len(layouts),
        "slide layout collection",
        0,
        1,
        "Re-fetch it with prs.slide_layouts.",
    )
    _assert_exact_stale(
        rpptx,
        lambda: layout.name,
        "slide layout",
        0,
        1,
        "Re-fetch it with prs.slide_layouts[6].",
    )
    shapes = slide.shapes
    shapes.add_textbox(0, 0, 100, 100)

    _assert_exact_stale(
        rpptx,
        lambda: slide.shapes,
        "slide",
        1,
        2,
        "Re-fetch it with prs.slides[0].",
    )
    _assert_exact_stale(
        rpptx,
        lambda: len(shapes),
        "shape collection",
        1,
        2,
        "Re-fetch it with prs.slides[0].shapes.",
    )


def test_stale_shape_text_and_table_paths_report_exact_recovery():
    import rpptx

    prs = rpptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    shape = prs.slides[0].shapes.add_textbox(0, 0, 100, 100)
    shape.text = "before"
    prs.slides[0].shapes.add_table(1, 1, 0, 0, 100, 100)

    slide = prs.slides[0]
    shapes = slide.shapes
    placeholders = slide.placeholders
    shape = shapes[0]
    frame = shape.text_frame
    paragraphs = frame.paragraphs
    paragraph = paragraphs[0]
    runs = paragraph.runs
    run = runs[0]
    font = paragraph.font
    table = shapes[1].table
    columns = table.columns
    column = columns[0]
    cell = table.cell(0, 0)
    prs.slides.add_slide(prs.slide_layouts[6])

    cases = (
        (lambda: slide.shapes, "slide", "Re-fetch it with prs.slides[0]."),
        (lambda: len(shapes), "shape collection", "Re-fetch it with prs.slides[0].shapes."),
        (lambda: placeholders[0], "placeholder collection", "Re-fetch it with prs.slides[0].placeholders."),
        (lambda: shape.text, "shape", "Re-fetch it with prs.slides[0].shapes[0]."),
        (lambda: frame.text, "text frame", "Re-fetch it with prs.slides[0].shapes[0].text_frame."),
        (lambda: len(paragraphs), "paragraph collection", "Re-fetch it with prs.slides[0].shapes[0].text_frame.paragraphs."),
        (lambda: paragraph.text, "paragraph", "Re-fetch it with prs.slides[0].shapes[0].text_frame.paragraphs[0]."),
        (lambda: len(runs), "run collection", "Re-fetch it with prs.slides[0].shapes[0].text_frame.paragraphs[0].runs."),
        (lambda: run.text, "run", "Re-fetch it with prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]."),
        (lambda: font.bold, "font", "Re-fetch it with prs.slides[0].shapes[0].text_frame.paragraphs[0].font."),
        (lambda: table.columns, "table", "Re-fetch it with prs.slides[0].shapes[1].table."),
        (lambda: len(columns), "column collection", "Re-fetch it with prs.slides[0].shapes[1].table.columns."),
        (lambda: column.width, "column", "Re-fetch it with prs.slides[0].shapes[1].table.columns[0]."),
        (lambda: cell.text, "cell", "Re-fetch it with prs.slides[0].shapes[1].table.cell(0, 0)."),
    )
    for operation, kind, recovery in cases:
        _assert_exact_stale(rpptx, operation, kind, 4, 5, recovery)


def test_whole_text_replacement_stales_descendant_handles_once():
    import rpptx

    prs = rpptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 100, 100)
    shape.text = "before"
    shape = prs.slides[0].shapes[-1]
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.runs[0]
    shape.text = "after"
    for operation in (
        lambda: shape.text,
        lambda: frame.text,
        lambda: paragraph.text,
        lambda: run.text,
    ):
        _assert_stale_after_exactly_one_bump(rpptx, operation)
    assert prs.slides[0].shapes[-1].text == "after"

    prs = rpptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
    frame.text = "before"
    frame = prs.slides[0].shapes[-1].text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.runs[0]
    frame.text = "after"
    for operation in (lambda: frame.text, lambda: paragraph.text, lambda: run.text):
        _assert_stale_after_exactly_one_bump(rpptx, operation)
    assert prs.slides[0].shapes[-1].text_frame.text == "after"

    prs = rpptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
    frame.text = "before"
    frame = prs.slides[0].shapes[-1].text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.runs[0]
    paragraph.text = "after"
    for operation in (lambda: paragraph.text, lambda: run.text):
        _assert_stale_after_exactly_one_bump(rpptx, operation)
    assert prs.slides[0].shapes[-1].text_frame.paragraphs[0].text == "after"


def test_nested_group_paths_report_exact_recovery(tmp_path):
    if importlib.util.find_spec("pptx") is None:
        pytest.skip("python-pptx oracle is installed only for the differential gate")

    import rpptx
    from pptx import Presentation as OraclePresentation

    source = tmp_path / "nested-groups.pptx"
    oracle = OraclePresentation()
    slide = oracle.slides.add_slide(oracle.slide_layouts[6])
    outer = slide.shapes.add_group_shape()
    inner = outer.shapes.add_group_shape()
    textbox = inner.shapes.add_textbox(0, 0, 100, 100)
    textbox.text = "nested"
    oracle.save(source)

    prs = rpptx.Presentation(source)
    inner = prs.slides[0].shapes[0].shapes[0]
    nested_shapes = inner.shapes
    shape = nested_shapes[0]
    frame = shape.text_frame
    paragraphs = frame.paragraphs
    paragraph = paragraphs[0]
    runs = paragraph.runs
    run = runs[0]
    font = paragraph.font
    prs.slides.add_slide(prs.slide_layouts[6])

    cases = (
        (lambda: len(nested_shapes), "shape collection", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes."),
        (lambda: shape.text, "shape", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0]."),
        (lambda: frame.text, "text frame", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame."),
        (lambda: len(paragraphs), "paragraph collection", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame.paragraphs."),
        (lambda: paragraph.text, "paragraph", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame.paragraphs[0]."),
        (lambda: len(runs), "run collection", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame.paragraphs[0].runs."),
        (lambda: run.text, "run", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame.paragraphs[0].runs[0]."),
        (lambda: font.bold, "font", "Re-fetch it with prs.slides[0].shapes[0].shapes[0].shapes[0].text_frame.paragraphs[0].font."),
    )
    for operation, kind, recovery in cases:
        _assert_exact_stale(rpptx, operation, kind, 0, 1, recovery)


def test_unit_constructors_truncate_fractional_values_toward_zero():
    import rpptx
    from rpptx.enum.shapes import MSO_SHAPE
    from rpptx.util import Inches, Length, Pt

    for constructor, factor in ((Inches, 914_400.0), (Pt, 12_700.0)):
        assert constructor(1.75 / factor) == 1
        assert constructor(-1.75 / factor) == -1
    assert Length(914_400).inches == 1.0
    assert Length(12_700).pt == 1.0
    assert Length(-1).emu == -1
    assert MSO_SHAPE.PENTAGON == 51
    assert MSO_SHAPE.CHEVRON == 52
    assert issubclass(rpptx.PackageError, rpptx.RpptxError)
    assert issubclass(rpptx.XmlError, rpptx.RpptxError)
    assert issubclass(rpptx.StaleElementError, rpptx.RpptxError)


def test_missing_package_raises_the_named_package_error(tmp_path):
    import rpptx

    with pytest.raises(rpptx.PackageError):
        rpptx.Presentation(tmp_path / "missing.pptx")


def _author_documented_decks(Presentation, Inches, Pt, MSO_SHAPE, root, image_path):
    root.mkdir()
    paths = {name: root / f"{name}.pptx" for name in (
        "hello", "bullet", "textbox", "picture", "shapes", "table"
    )}

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle = prs.slides[0].placeholders[1]
    subtitle.text = "python-pptx was here!"
    prs.save(paths["hello"])

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    shapes.title.text = "Adding a Bullet Slide"
    shapes = prs.slides[0].shapes
    frame = shapes.placeholders[1].text_frame
    frame.text = "Find the bullet slide layout"
    frame = prs.slides[0].shapes.placeholders[1].text_frame
    paragraph = frame.add_paragraph()
    paragraph.text = "Use _TextFrame.text for first bullet"
    paragraph = prs.slides[0].shapes.placeholders[1].text_frame.paragraphs[1]
    paragraph.level = 1
    frame = prs.slides[0].shapes.placeholders[1].text_frame
    paragraph = frame.add_paragraph()
    paragraph.text = "Use _TextFrame.add_paragraph() for subsequent bullets"
    paragraph = prs.slides[0].shapes.placeholders[1].text_frame.paragraphs[2]
    paragraph.level = 2
    prs.save(paths["bullet"])

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)).text_frame
    frame.text = "This is text inside a textbox"
    frame = prs.slides[0].shapes[-1].text_frame
    paragraph = frame.add_paragraph()
    paragraph.text = "This is a second paragraph that's bold"
    paragraph = prs.slides[0].shapes[-1].text_frame.paragraphs[1]
    paragraph.font.bold = True
    frame = prs.slides[0].shapes[-1].text_frame
    paragraph = frame.add_paragraph()
    paragraph.text = "This is a third paragraph that's big"
    paragraph = prs.slides[0].shapes[-1].text_frame.paragraphs[2]
    paragraph.font.size = Pt(40)
    prs.save(paths["textbox"])

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    slide = prs.slides[0]
    slide.shapes.add_picture(
        str(image_path), Inches(5), Inches(1), height=Inches(5.5)
    )
    prs.save(paths["picture"])

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    shapes.title.text = "Adding an AutoShape"
    shapes = prs.slides[0].shapes
    left = Inches(0.93)
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)
    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = "Step 1"
    left = left + width - Inches(0.4)
    width = Inches(2.0)
    for number in range(2, 6):
        shapes = prs.slides[0].shapes
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = f"Step {number}"
        left = left + width - Inches(0.4)
    prs.save(paths["shapes"])

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    shapes.title.text = "Adding a Table"
    shapes = prs.slides[0].shapes
    table = shapes.add_table(
        2, 2, Inches(2), Inches(2), Inches(6), Inches(0.8)
    ).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4)
    for row, values in enumerate((("Foo", "Bar"), ("Baz", "Qux"))):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    prs.save(paths["table"])
    return paths


def _normalized_table(table):
    columns = len(table.columns)
    rows = []
    row = 0
    while True:
        try:
            rows.append(tuple(table.cell(row, column).text for column in range(columns)))
        except IndexError:
            break
        row += 1
    return tuple(int(column.width) for column in table.columns), tuple(rows)


def _normalized_presentation(prs):
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            paragraphs = ()
            if shape.has_text_frame:
                paragraphs = tuple(
                    (
                        paragraph.text,
                        paragraph.level,
                        paragraph.font.bold,
                        int(paragraph.font.size) if paragraph.font.size is not None else None,
                        tuple(run.text for run in paragraph.runs),
                    )
                    for paragraph in shape.text_frame.paragraphs
                )
            shapes.append(
                (
                    shape.has_text_frame,
                    paragraphs,
                    shape.has_table,
                    _normalized_table(shape.table) if shape.has_table else None,
                )
            )
        slides.append(tuple(shapes))
    return tuple(slides)


def _normalized_text_extraction(prs):
    return tuple(
        run.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )


def _normalized_documented_records(Presentation, paths):
    records = {
        name: _normalized_presentation(Presentation(path))
        for name, path in paths.items()
    }
    records["extract"] = _normalized_text_extraction(Presentation(paths["table"]))
    return records


def _oracle_writer_contract(Presentation, paths):
    def auto_shape_type(shape):
        if int(shape.shape_type) != 1:
            return None
        try:
            return int(shape.auto_shape_type)
        except (AttributeError, TypeError, ValueError):
            return None

    records = {}
    for name, path in paths.items():
        prs = Presentation(path)
        records[name] = tuple(
            tuple(
                (
                    int(shape.shape_type),
                    (int(shape.left), int(shape.top), int(shape.width), int(shape.height)),
                    auto_shape_type(shape),
                    int(shape.placeholder_format.idx) if shape.is_placeholder else None,
                )
                for shape in slide.shapes
            )
            for slide in prs.slides
        )
    return records


def test_pinned_python_pptx_bidirectional_seven_example_records(tmp_path):
    if importlib.util.find_spec("pptx") is None:
        pytest.skip("python-pptx oracle is installed only for the differential gate")
    assert importlib.metadata.version("python-pptx") == "1.0.2"

    import rpptx
    from rpptx.enum.shapes import MSO_SHAPE as RpptxShape
    from rpptx.util import Inches as RpptxInches
    from rpptx.util import Pt as RpptxPt
    from pptx import Presentation as OraclePresentation
    from pptx.enum.shapes import MSO_SHAPE as OracleShape
    from pptx.util import Inches as OracleInches
    from pptx.util import Pt as OraclePt

    image_path = tmp_path / "monty-truth.png"
    _write_tiny_png(image_path)
    authored = (
        _author_documented_decks(
            rpptx.Presentation,
            RpptxInches,
            RpptxPt,
            RpptxShape,
            tmp_path / "rpptx-authored",
            image_path,
        ),
        _author_documented_decks(
            OraclePresentation,
            OracleInches,
            OraclePt,
            OracleShape,
            tmp_path / "python-pptx-authored",
            image_path,
        ),
    )
    normalized = []
    writer_contracts = []
    for paths in authored:
        rpptx_records = _normalized_documented_records(rpptx.Presentation, paths)
        oracle_records = _normalized_documented_records(OraclePresentation, paths)
        assert rpptx_records == oracle_records
        normalized.append(rpptx_records)
        writer_contracts.append(_oracle_writer_contract(OraclePresentation, paths))
    assert normalized[0] == normalized[1]
    assert writer_contracts[0] == writer_contracts[1]
    assert normalized[0]["extract"] == ("Adding a Table",)
